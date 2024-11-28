import os.path

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE

from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt

from xpinyin import Pinyin

def getInches(data_cm) -> Inches:
    return data_cm / 2.54

def addTitle(slide_obj, text_data):
    left = Inches(getInches(0))
    top = Inches(getInches(0))
    width = Inches(getInches(18.39))
    height = Inches(getInches(1.8))
    title_textbox = slide_obj.shapes.add_textbox(left, top, width, height)
    text_frame = title_textbox.text_frame
    text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT # 自适应大小
    p = text_frame.paragraphs[0]
    p.text = text_data  # 添加的文本内容
    p.font.bold = True  # 设置字体为粗体
    p.font.size = Inches(0.5)  # 设置字体大小 1英寸=72磅
    p.font.name = '华文新魏'  # 设置字体名称
    p.font.color.rgb = RGBColor(192, 0, 0)  # 设置字体颜色为红色

def addContextSub(p, head_text, body_text):
    run = p.add_run()
    run.text = head_text
    font = run.font
    font.name = '华文新魏'
    font.size = Pt(28)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    font.color.rgb = RGBColor(192, 0, 0)

    run = p.add_run()
    run.text = body_text
    font = run.font
    font.name = '华文新魏'
    font.size = Pt(28)
    font.bold = False
    font.italic = None  # cause value to be inherited from theme
    font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    font.color.rgb = RGBColor(0, 0, 255)


def addContext(slide_obj, text_dingwei, text_jiepou, text_zhuzhi, text_cifa):

    left = Inches(getInches(0))
    top = Inches(getInches(1.8))
    width = Inches(getInches(13))
    height = Inches(getInches(15))
    title_textbox = slide_obj.shapes.add_textbox(left, top, width, height)
    text_frame = title_textbox.text_frame
    text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT  # 自适应大小

    p0 = text_frame.paragraphs[0]
    addContextSub(p0, "定位：", text_dingwei)

    p1 = text_frame.add_paragraph()
    addContextSub(p1, "解剖：", text_jiepou)

    p2 = text_frame.add_paragraph()
    addContextSub(p2, "主治：", text_zhuzhi)

    p3 = text_frame.add_paragraph()
    addContextSub(p3, "刺法：", text_cifa)

    p4 = text_frame.add_paragraph()
    addContextSub(p4, "灸法：", "TODO")


def addOneFrame(ppt_obj, meta_file_path, title_text):
    # 添加一张空白幻灯片
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # 添加标题
    addTitle(slide, title_text)



    # 添加备注
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    lines = None
    with open(meta_file_path, "r", encoding='utf-8') as file_handle:
        lines = file_handle.readlines()
        file_handle.seek(0)
        tmp_lines = file_handle.read()
        text_frame.text = tmp_lines


    text_dingwei = "TODO"
    text_jiepou = "TODO"
    text_zhuzhi = "TODO"
    text_cifa = "TODO"
    for line in lines:
        if line[1:3] == "定位":
            text_dingwei = line[4:-1]
        if line[1:3] == "解剖":
            text_jiepou = line[4:-1]
        if line[1:3] == "主治":
            text_zhuzhi = line[4:-1]
        if line[1:4] == "刺灸法":
            text_cifa = line[5:-1]

    print(">>> 定位 {}".format(text_dingwei))
    print(">>> 解剖 {}".format(text_jiepou))
    print(">>> 主治 {}".format(text_zhuzhi))
    print(">>> 针刺 {}".format(text_cifa))
    # 添加内容
    addContext(slide, text_dingwei, text_jiepou, text_zhuzhi, text_cifa)


if __name__ == "__main__":
    prs = Presentation()
    ST_ACUPOINT = ["承泣穴", "四白穴", "巨髎穴", "地仓穴", "大迎穴", "颊车穴", "下关穴", "头维穴", "人迎穴", "水突穴",
                   "气舍穴", "缺盆穴", "气户穴", "库房穴", "屋翳穴", "膺窗穴",  "乳中穴", "乳根穴", "不容穴", "承满穴",
                   "梁门穴", "关门穴", "太乙穴", "滑肉门穴", "天枢穴", "外陵穴", "大巨穴", "水道穴", "归来穴",
                   "气冲穴", "髀关穴", "伏兔穴", "阴市穴", "梁丘穴", "犊鼻穴", "足三里穴", "上巨虚穴", "条口穴",
                   "下巨虚穴", "丰隆穴", "解溪穴", "冲阳穴", "陷谷穴", "内庭穴", "厉兑穴"]

    SP_ACUPOINT = ["隐白穴", "大都穴", "太白穴", "公孙穴", "商丘穴", "三阴交穴", "漏谷穴",
                   "地机穴", "阴陵泉穴", "血海穴", "箕门穴", "冲门穴", "府舍穴", "腹结穴",
                   "大横穴", "腹哀穴", "食窦穴", "天溪穴", "胸乡穴", "周荣穴", "大包穴"]

    HT_ACUPOINT = ["极泉穴", "青灵穴", "少海穴", "灵道穴", "通里穴", "阴郄穴", "神门穴", "少府穴", "少冲穴"]
    SI_ACUPOINT = ["少泽穴", "前谷穴", "后溪穴", "腕骨穴", "阳谷穴", "养老穴", "支正穴", "小海穴", "肩贞穴", "臑俞穴",
                   "天宗穴", "秉风穴", "曲垣穴", "肩外俞穴", "肩中俞穴", "天窗穴", "天容穴", "颧髎穴", "听宫穴"]

    sz = len(SI_ACUPOINT)
    for index in range(len(SI_ACUPOINT)):
        acupoint_name = SI_ACUPOINT[index]
        ele_meta_file = "data" + os.sep + "手太阳小肠经" + os.sep + acupoint_name + ".txt"
        pinyin_obj = Pinyin()
        result = pinyin_obj.get_pinyin(acupoint_name)
        result_1 = result.split('-')
        R = ""
        for py_ele in result_1:
            R = R + py_ele[0].upper() + py_ele[1:]

        title_text = "{}-{} {} ({}, ST{})".format(len(SI_ACUPOINT), index+1, acupoint_name, R, index+1)
        addOneFrame(prs, ele_meta_file, title_text)



    prs.save('手太阳小肠经.pptx')